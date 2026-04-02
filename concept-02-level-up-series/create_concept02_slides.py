#!/usr/bin/env python3
"""Create Concept 02 — Level Up Series updated slides."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Design System Colours ──
INK       = RGBColor(0x26, 0x25, 0x22)
CHARCOAL  = RGBColor(0x3D, 0x39, 0x35)
STORM     = RGBColor(0x7B, 0x78, 0x74)
CLOUD     = RGBColor(0xB2, 0xB0, 0xAD)
LINEN     = RGBColor(0xF5, 0xF4, 0xF1)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
CORAL     = RGBColor(0xFF, 0x6B, 0x6B)
CORAL_DK  = RGBColor(0xD9, 0x40, 0x40)
YELLOW    = RGBColor(0xFF, 0xCE, 0x49)
YELLOW_DK = RGBColor(0xC5, 0x96, 0x0C)
YELLOW_LT = RGBColor(0xFC, 0xF3, 0xDB)
GREEN     = RGBColor(0x45, 0xBE, 0x93)
GREEN_DK  = RGBColor(0x2D, 0x8A, 0x6A)
GREEN_LT  = RGBColor(0xCE, 0xEF, 0xE4)
VIOLET    = RGBColor(0x84, 0x5E, 0xF7)
VIOLET_DK = RGBColor(0x5B, 0x3D, 0xB8)
VIOLET_LT = RGBColor(0xE5, 0xDB, 0xFF)
BLUE      = RGBColor(0x1C, 0x83, 0xE1)
BLUE_DK   = RGBColor(0x0D, 0x4F, 0x8A)
BLUE_LT   = RGBColor(0xE9, 0xF3, 0xFD)
AMBER     = RGBColor(0xE8, 0x9D, 0x2D)
AMBER_LT  = RGBColor(0xFD, 0xF0, 0xD0)
CORAL_LT  = RGBColor(0xFF, 0xE0, 0xE0)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = 13.333
H = 7.5

# ── Helpers ──
def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def rrect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def circle(slide, x, y, sz, color):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(sz), Inches(sz))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def text(slide, txt, x, y, w, h, sz=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font='Calibri'):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = txt
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold
    p.font.name = font; p.alignment = align
    return tb

def bullets(slide, items, x, y, w, h, sz=11, color=STORM, spacing=2):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(sz); p.font.color.rgb = color
        p.font.name = 'Calibri'; p.space_after = Pt(spacing)
    return tb

def label(slide, txt, x, y, color=VIOLET, w=3):
    rrect(slide, x, y, w, 0.32, color)
    text(slide, txt, x, y, w, 0.32, sz=9, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, LINEN)

# Header bar
rect(s, 0, 0, W, 3.2, CORAL)

text(s, "CONCEPT 02 — UPDATED", 1.0, 0.5, 6, 0.3, sz=11, color=WHITE, bold=True)
text(s, "Level Up Series", 1.0, 1.0, 8, 1.0, sz=52, color=WHITE, bold=True)
text(s, "What changed: Wellness Board (C1) + Made for Parents (C4)\nnow live inside the Level Up Series box.", 1.0, 2.1, 8, 0.7, sz=14, color=WHITE)

# Bottom text
text(s, "HPB YPS  |  March 2026  |  Updated Concept Slides", 1.0, 6.8, 6, 0.4, sz=11, color=STORM)


# ═══════════════════════════════════════
# SLIDE 2: WHAT THE BOX NOW COMBINES
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, LINEN)

label(s, "WHAT CHANGED", 0.8, 0.4, CORAL)
text(s, "Three Concepts, One Box", 0.8, 0.85, 10, 0.6, sz=32, color=INK, bold=True)
text(s, "Level Up Series now combines C1 (Wellness Board), C3 (Level Up tools), and C4 (Made for Parents) into a single physical kit.", 0.8, 1.45, 10, 0.4, sz=13, color=STORM)

# Three source cards
sources = [
    ("C1", "Wellness Board", "Physical magnetic habit boards.\nYouth-owned, one well at a time.\nNow includes 5 boards + Exam Season.", GREEN, GREEN_LT, GREEN_DK),
    ("C3", "Level Up Series Tools", "Start the Talk (32 cards),\nStart Simple (16 challenges + scoreboard),\nHow Are You? QR, Quick Start Guide.", CORAL, CORAL_LT, CORAL_DK),
    ("C4", "Made for Parents", "Real parent videos, situation-matched\ntips, weekly WhatsApp nudges.\nOne QR card in the box. Not a lecture.", VIOLET, VIOLET_LT, VIOLET_DK),
]
for i, (tag, title, desc, accent, bg, dark) in enumerate(sources):
    cx = 0.8 + i * 4.0
    rrect(s, cx, 2.1, 3.75, 3.6, WHITE)
    rect(s, cx, 2.1, 3.75, 0.07, accent)
    # Tag
    circle(s, cx + 0.15, 2.3, 0.55, accent)
    text(s, tag, cx + 0.15, 2.3, 0.55, 0.55, sz=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # Originally label
    text(s, f"Originally {tag}", cx + 0.85, 2.35, 2, 0.3, sz=9, color=accent, bold=True)
    # Title
    text(s, title, cx + 0.2, 2.85, 3.3, 0.45, sz=16, color=INK, bold=True)
    # Desc
    text(s, desc, cx + 0.2, 3.35, 3.3, 1.8, sz=11, color=STORM)

# Arrow section
rrect(s, 0.8, 6.0, 11.7, 0.9, INK)
text(s, "C1 + C3 + C4  =  Level Up Series (one box)          C2 = Digital Challenge / Challenge Accepted (separate, stays as-is)", 1.1, 6.15, 11, 0.55, sz=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════
# SLIDE 3: WHAT'S NOW IN THE BOX
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, LINEN)

label(s, "UPDATED BOX CONTENTS", 0.8, 0.4, CORAL)
text(s, "What's Inside — 6 Components", 0.8, 0.85, 10, 0.6, sz=32, color=INK, bold=True)

components = [
    ("1", "Start the Talk", "32 conversation cards\n'Never Have I Ever' style", YELLOW, YELLOW_LT),
    ("2", "Start Simple", "16 challenge cards +\nmagnetic scoreboard", CORAL, CORAL_LT),
    ("3", "How Are You?", "QR mood check-in card\nDigital bridge", GREEN, GREEN_LT),
    ("4", "Wellness Boards", "5 magnetic boards\nSleep · Eat · Move · Mind · Exam", BLUE, BLUE_LT),
    ("5", "Parent Zone", "Made for Parents QR card\nVideos, tips, WhatsApp nudges", VIOLET, VIOLET_LT),
    ("6", "Quick Start Guide", "Folded instruction card\n'Start here'", INK, RGBColor(0xF0, 0xEF, 0xEC)),
]
for i, (num, name, desc, accent, bg) in enumerate(components):
    row = i // 3
    col = i % 3
    cx = 0.8 + col * 4.0
    cy = 1.7 + row * 2.7

    rrect(s, cx, cy, 3.75, 2.4, bg)
    rect(s, cx, cy, 0.07, 2.4, accent)
    circle(s, cx + 0.2, cy + 0.2, 0.5, accent)
    text(s, num, cx + 0.2, cy + 0.2, 0.5, 0.5, sz=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    text(s, name, cx + 0.85, cy + 0.2, 2.6, 0.4, sz=15, color=INK, bold=True)
    text(s, desc, cx + 0.85, cy + 0.65, 2.6, 1.2, sz=11, color=STORM)

    # NEW badge on items 4 and 5
    if i in [3, 4]:
        rrect(s, cx + 2.8, cy + 0.15, 0.7, 0.25, CORAL)
        text(s, "NEW", cx + 2.8, cy + 0.15, 0.7, 0.25, sz=8, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════
# SLIDE 4: WELLNESS BOARDS — 5TH BOARD
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, LINEN)

label(s, "NEW — WELLNESS BOARDS", 0.8, 0.4, BLUE)
text(s, "5 Boards — Now Includes Exam Season", 0.8, 0.85, 10, 0.6, sz=32, color=INK, bold=True)
text(s, "Youth-owned. Lives in their space. Parents only come in if invited. During exam season, swap in the 5th board.", 0.8, 1.45, 10, 0.4, sz=13, color=STORM)

boards = [
    ("Better Sleep", BLUE, BLUE_LT, "Phone down by 9:30pm\n8 hours sleep target\nDim lights after 8pm\nSame bedtime 7 nights"),
    ("Eat Better", GREEN, GREEN_LT, "Eat breakfast every day\nDrink 8 glasses of water\nCook one meal together\nTry one new fruit/veg"),
    ("Move Your Body", CORAL, CORAL_LT, "10-min walk after dinner\nStretch 5 min before bed\nDance to one song a day\nTake the stairs"),
    ("Mind & Focus", VIOLET, VIOLET_LT, "No social media 1hr after waking\n5-min breathing before bed\n1 gratitude entry daily\nDinner without phones"),
    ("Exam Season", AMBER, AMBER_LT, "Study break every 45 min\nNo caffeine after 3pm\nProtect sleep — 7hrs minimum\nOne meal away from desk\n10-min walk between subjects"),
]
for i, (name, accent, bg, habits) in enumerate(boards):
    cx = 0.8 + i * 2.45
    rrect(s, cx, 2.1, 2.2, 4.5, WHITE)
    rect(s, cx, 2.1, 2.2, 0.08, accent)

    # Board name
    text(s, name, cx + 0.15, 2.3, 1.9, 0.35, sz=13, color=INK, bold=True)

    # Habits list
    habit_list = habits.split("\n")
    for j, habit in enumerate(habit_list):
        circle(s, cx + 0.15, 2.8 + j * 0.45, 0.12, accent)
        text(s, habit, cx + 0.35, 2.75 + j * 0.45, 1.7, 0.4, sz=9, color=STORM)

    # NEW badge on Exam Season
    if i == 4:
        rrect(s, cx + 1.4, 2.2, 0.6, 0.22, CORAL)
        text(s, "NEW", cx + 1.4, 2.2, 0.6, 0.22, sz=7, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Bottom note
rrect(s, 0.8, 6.8, 11.7, 0.5, AMBER_LT)
text(s, "Exam Season board: During revision, scale down to 1 board. Sleep is the one worth protecting — everything else improves when sleep is right.", 1.0, 6.85, 11, 0.35, sz=10, color=INK, bold=True)


# ═══════════════════════════════════════
# SLIDE 5: MADE FOR PARENTS — QR CARD
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, LINEN)

label(s, "NEW — MADE FOR PARENTS", 0.8, 0.4, VIOLET)
text(s, "One QR Card. Real Parent Content.", 0.8, 0.85, 10, 0.6, sz=32, color=INK, bold=True)
text(s, "Originally Concept 4. Now lives inside the Level Up Series box as a single QR card. Parents scan it to access situation-matched content.", 0.8, 1.45, 10, 0.4, sz=13, color=STORM)

# Left: What parents get
rrect(s, 0.8, 2.1, 5.6, 4.5, WHITE)
rect(s, 0.8, 2.1, 5.6, 0.07, VIOLET)
text(s, "What Parents Get", 1.0, 2.3, 4, 0.4, sz=16, color=INK, bold=True)

features = [
    ("Real parent videos", "Not scripted. A parent talking from their kitchen, not a studio. Raw, honest, practical."),
    ("Situation-matched tips", "Bedtime Battles, Silent Dinners, The Screen Fight, Exam Season — pick what you're dealing with now."),
    ("Transition-stage content", "Tips adjust by age: Sec 1 (still open), Sec 2-3 (pulling away), Sec 4+ (self-managing)."),
    ("Weekly WhatsApp nudge", "One relevant tip per week. Linked to a situation matching their child's age and stage."),
    ("'Did this help?' feedback", "Simple tap after each tip so content improves over time."),
]
for i, (title, desc) in enumerate(features):
    cy = 2.85 + i * 0.7
    circle(s, 1.0, cy, 0.12, VIOLET)
    text(s, title, 1.25, cy - 0.05, 2, 0.3, sz=11, color=INK, bold=True)
    text(s, desc, 1.25, cy + 0.2, 4.8, 0.4, sz=9, color=STORM)

# Right: Phone mockup representation
rrect(s, 6.8, 2.1, 5.6, 4.5, VIOLET_LT)
text(s, "PHONE SCREEN MOCKUP", 7.0, 2.3, 4, 0.3, sz=9, color=VIOLET, bold=True)

# Mockup cards
situations = [
    ("This week", "Bedtime Battles", "The phone won't go down"),
    ("", "Silent Dinners", "When the table goes quiet"),
    ("", "The Screen Fight", "Every convo becomes about screen time"),
    ("", "Exam Season", "Everything healthy goes out the window"),
]
for i, (tag, title, desc) in enumerate(situations):
    cy = 2.8 + i * 0.85
    rrect(s, 7.2, cy, 4.8, 0.7, WHITE)
    if tag:
        text(s, tag, 7.4, cy + 0.05, 2, 0.2, sz=8, color=VIOLET, bold=True)
    text(s, title, 7.4, cy + 0.2, 3, 0.25, sz=12, color=INK, bold=True)
    text(s, desc, 7.4, cy + 0.42, 3, 0.2, sz=9, color=STORM)

# Bottom
text(s, "Powered by Parenting for Wellness  |  levelup.sg/parents", 7.2, 6.2, 4, 0.3, sz=9, color=VIOLET)


# ═══════════════════════════════════════
# SLIDE 6: TRANSITION FRAMEWORK
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, LINEN)

label(s, "HOW THE BOX ADAPTS", 0.8, 0.4, INK, w=3.5)
text(s, "Across Three Transitions", 0.8, 0.85, 10, 0.6, sz=32, color=INK, bold=True)
text(s, "The box arrives at transition points — not after. Each tool plays a different role depending on where the youth is.", 0.8, 1.45, 10, 0.4, sz=13, color=STORM)

transitions = [
    ("T1", "P6 → Sec 1 (12–13)", "Parent-driven", CORAL, [
        "Start the Talk + Start Simple land — relationship still open",
        "Silly tone makes it safe enough to start",
        "Wellness Board in youth's room — parent sees, youth owns",
        "Parent Zone QR helps parents adjust from enforcing to supporting",
    ]),
    ("T2", "Sec 2 → Sec 3 (14–15)", "Peer-driven", BLUE, [
        "Peer accountability replaces parent-led tools",
        "Swap Cards essential — answering as each other is less exposing",
        "Challenge cards shift to peer pairs, not parent-youth",
        "Digital Challenge (Challenge Accepted) becomes primary tool",
        "Board becomes fully youth-owned; parent visibility drops",
    ]),
    ("T3", "Sec 4 → JC/IHL (16+)", "Self-managing", GREEN, [
        "Youth is self-managing — card games no longer needed",
        "Wellness Board + How Are You? QR become core tools",
        "Exam Season board most relevant — O/A-level pressure",
        "Parent Zone shifts to 'how to step back without disconnecting'",
    ]),
]
for i, (tag, period, role, accent, items) in enumerate(transitions):
    cx = 0.8 + i * 4.0
    rrect(s, cx, 2.0, 3.75, 4.8, WHITE)
    rect(s, cx, 2.0, 3.75, 0.08, accent)

    # Tag + period
    circle(s, cx + 0.15, 2.25, 0.5, accent)
    text(s, tag, cx + 0.15, 2.25, 0.5, 0.5, sz=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    text(s, period, cx + 0.8, 2.3, 2.7, 0.3, sz=11, color=accent, bold=True)
    text(s, role, cx + 0.8, 2.6, 2.7, 0.3, sz=16, color=INK, bold=True)

    # Bullets
    for j, item in enumerate(items):
        cy = 3.1 + j * 0.45
        circle(s, cx + 0.2, cy + 0.05, 0.08, accent)
        text(s, item, cx + 0.4, cy - 0.02, 3.1, 0.4, sz=9, color=STORM)

# Bottom principle
rrect(s, 0.8, 7.0, 11.7, 0.35, VIOLET_LT)
text(s, "Parent involvement doesn't disappear — it changes form. From structure-setter (T1) to silent supporter (T2) to available-if-needed (T3).", 1.0, 7.02, 11, 0.3, sz=10, color=INK, bold=True)


# ═══════════════════════════════════════
# SLIDE 7: HOW C2 + LUS CONNECT
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, LINEN)

label(s, "TWO CONCEPTS, ONE SYSTEM", 0.8, 0.4, INK, w=4)
text(s, "How Level Up Series & Digital Challenge Connect", 0.8, 0.85, 10, 0.6, sz=28, color=INK, bold=True)

# Left: Level Up Series
rrect(s, 0.8, 1.7, 5.6, 4.8, WHITE)
rect(s, 0.8, 1.7, 5.6, 0.08, CORAL)
text(s, "Level Up Series", 1.0, 1.95, 4, 0.5, sz=20, color=INK, bold=True)
text(s, "Physical box  |  Roadshow distribution  |  C1 + C3 + C4", 1.0, 2.4, 5, 0.3, sz=10, color=CORAL, bold=True)

lus_items = [
    "Start the Talk — 32 conversation cards",
    "Start Simple — 16 challenge cards + scoreboard",
    "Wellness Boards — 5 magnetic boards (incl. Exam Season)",
    "How Are You? — QR mood check-in",
    "Parent Zone — Made for Parents QR card",
    "Quick Start Guide — folded instruction card",
]
bullets(s, lus_items, 1.0, 2.9, 5, 3.0, sz=11, color=STORM, spacing=6)

# Right: Digital Challenge
rrect(s, 6.8, 1.7, 5.6, 4.8, WHITE)
rect(s, 6.8, 1.7, 5.6, 0.08, VIOLET)
text(s, "Digital Challenge", 7.0, 1.95, 4, 0.5, sz=20, color=INK, bold=True)
text(s, "Web-based  |  QR or friend invite entry  |  C2 (unchanged)", 7.0, 2.4, 5, 0.3, sz=10, color=VIOLET, bold=True)

dc_items = [
    "Challenge Accepted — Duolingo-style gamified streaks",
    "5 packs: Glow Up, Energy Era, Cut Season, Bulk Up, Recharge",
    "Buddy accountability — shared streaks",
    "XP, leaderboards, celebrations",
    "14-day challenge cycle",
    "Light-mode, web-based (not native app)",
]
bullets(s, dc_items, 7.0, 2.9, 5, 3.0, sz=11, color=STORM, spacing=6)

# Connection arrow section
rrect(s, 0.8, 6.7, 11.7, 0.6, INK)
text(s, "Connection: Level Up Series QR cards → Digital Challenge entry point  |  Challenge cards bridge physical → digital  |  Both share the 6 Wells framework", 1.0, 6.78, 11, 0.4, sz=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════
# SLIDE 8: ROADSHOW — TWO EXPERIENCES
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, LINEN)

label(s, "ROADSHOW PLAN", 0.8, 0.4, CORAL)
text(s, "Two Roadshows, Two Audiences", 0.8, 0.85, 10, 0.6, sz=32, color=INK, bold=True)

# Roadshow 1
rrect(s, 0.8, 1.7, 5.6, 5.0, WHITE)
rect(s, 0.8, 1.7, 5.6, 0.08, YELLOW)
text(s, "The Family Experience", 1.0, 1.95, 4, 0.5, sz=18, color=INK, bold=True)
text(s, "ROADSHOW 1", 1.0, 2.4, 3, 0.25, sz=9, color=YELLOW_DK, bold=True)

r1_items = [
    "Who: Parents + Primary/Early secondary youth",
    "Where: PFW events, school orientation, HPB sites",
    "",
    "Level Up Series box — open & try Start the Talk on the spot",
    "Wellness Board demo — youth pick what goes on it",
    "Made for Parents corner — one video on screen, tip card to take home",
    "Challenge card — scan to start Digital Challenge (C2)",
]
bullets(s, r1_items, 1.0, 2.8, 5, 3.5, sz=10, color=STORM, spacing=3)

# Roadshow 2
rrect(s, 6.8, 1.7, 5.6, 5.0, WHITE)
rect(s, 6.8, 1.7, 5.6, 0.08, VIOLET)
text(s, "The Youth Activation", 7.0, 1.95, 4, 0.5, sz=18, color=INK, bold=True)
text(s, "ROADSHOW 2", 7.0, 2.4, 3, 0.25, sz=9, color=VIOLET_DK, bold=True)

r2_items = [
    "Who: Youth only (Sec 2+)",
    "Where: Malls, school events, community spaces",
    "",
    "Digital Challenge leaderboard on screen",
    "Challenge cards handed out by youth ambassadors",
    "QR → join a streak with a friend on the spot",
    "Wellness Board 'build your own' station",
    "No parent presence — youth-owned space",
]
bullets(s, r2_items, 7.0, 2.8, 5, 3.5, sz=10, color=STORM, spacing=3)


# ═══════════════════════════════════════
# SLIDE 9: SUMMARY
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, LINEN)

# Header bar
rect(s, 0, 0, W, 3.0, INK)

text(s, "Summary of Changes", 1.0, 0.8, 8, 0.8, sz=36, color=WHITE, bold=True)
text(s, "Level Up Series — Updated Concept 02", 1.0, 1.5, 6, 0.4, sz=14, color=STORM)

# Changes grid
changes = [
    ("Wellness Board (C1)", "Now inside the box", "5 boards (Sleep, Eat, Move, Mind, Exam Season). Youth-owned. Parents only if invited.", GREEN),
    ("Made for Parents (C4)", "Now a QR card in the box", "Real videos, situation-matched tips, WhatsApp nudges. One card, full digital playbook.", VIOLET),
    ("Exam Season Board", "New 5th board", "Study breaks, sleep protection, no caffeine after 3pm. Scales down during revision.", AMBER),
    ("Transition Framework", "New section", "Box adapts across T1 (parent-driven), T2 (peer-driven), T3 (self-managing).", BLUE),
    ("Digital Challenge (C2)", "No changes", "Challenge Accepted stays as separate web-based concept. Connected via QR.", CORAL),
]

for i, (title, status, desc, accent) in enumerate(changes):
    row = i // 3
    col = i % 3
    cx = 0.8 + col * 4.0
    cy = 3.5 + row * 1.8

    rrect(s, cx, cy, 3.75, 1.55, WHITE)
    rect(s, cx, cy, 0.07, 1.55, accent)
    text(s, title, cx + 0.2, cy + 0.12, 3.3, 0.3, sz=13, color=INK, bold=True)
    rrect(s, cx + 0.2, cy + 0.45, 1.8, 0.22, accent)
    text(s, status, cx + 0.2, cy + 0.45, 1.8, 0.22, sz=8, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    text(s, desc, cx + 0.2, cy + 0.75, 3.3, 0.7, sz=9, color=STORM)


# ── SAVE ──
out_path = "/Users/alveriacher/Library/CloudStorage/GoogleDrive-alveria.cher@aleph-labs.com/My Drive/HPB YPS/All info/YPSSSSS_Files/Concept 02-  Level Up Series/Concept_02_Level_Up_Series_Updated.pptx"
prs.save(out_path)
print(f"Saved to {out_path}")
print(f"Total slides: {len(prs.slides)}")
