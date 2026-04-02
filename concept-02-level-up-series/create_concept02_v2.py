#!/usr/bin/env python3
"""Concept 02 — Level Up Series (4 core + 2 add-ons) slides matching iterated deck structure."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colours (matching deck style — clean white bg) ──
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLACK     = RGBColor(0x1A, 0x1A, 0x1A)
DARK      = RGBColor(0x2D, 0x2D, 0x2D)
GREY      = RGBColor(0x6B, 0x6B, 0x6B)
GREY_LT   = RGBColor(0x9C, 0xA3, 0xAF)
GREY_BG   = RGBColor(0xF3, 0xF4, 0xF6)
CORAL     = RGBColor(0xFF, 0x6B, 0x6B)
YELLOW    = RGBColor(0xFF, 0xCE, 0x49)
GREEN     = RGBColor(0x45, 0xBE, 0x93)
BLUE      = RGBColor(0x1C, 0x83, 0xE1)
VIOLET    = RGBColor(0x84, 0x5E, 0xF7)
AMBER     = RGBColor(0xE8, 0x9D, 0x2D)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = 13.333; H = 7.5

def bg(slide, color=WHITE):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def rect(s, x, y, w, h, c):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()
    return sh

def rrect(s, x, y, w, h, c):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()
    return sh

def circ(s, x, y, sz, c):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(sz), Inches(sz))
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()
    return sh

def txt(s, t, x, y, w, h, sz=14, c=BLACK, bold=False, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = t; p.font.size = Pt(sz); p.font.color.rgb = c
    p.font.bold = bold; p.font.name = 'Calibri'; p.alignment = align
    return tb

def bullets(s, items, x, y, w, h, sz=12, c=GREY, sp=4):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(sz); p.font.color.rgb = c
        p.font.name = 'Calibri'; p.space_after = Pt(sp)
    return tb

def tag(s, t, x, y, w=2.5, c=VIOLET):
    rrect(s, x, y, w, 0.28, c)
    txt(s, t, x, y, w, 0.28, sz=8, c=WHITE, bold=True, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# SLIDE 1: CONCEPT COVER
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)

# Large concept number
txt(s, "Concept 02", 0.8, 0.5, 4, 0.5, sz=14, c=GREY_LT)

# Concept title
txt(s, "Level Up\nSeries", 0.8, 1.5, 6, 2.5, sz=56, c=BLACK, bold=True)

# Tagline
txt(s, "Build habits at the moments\nthat matter — together.", 0.8, 4.3, 6, 1.0, sz=18, c=GREY)

# Note at bottom
txt(s, "Base box: 4 core tools  |  Optional add-ons: Wellness Boards (C1) + Parent Zone (C4)", 0.8, 6.2, 6, 0.4, sz=11, c=GREY_LT)
txt(s, "Image is for illustration purposes only\nRefer to appendix F for full details", 0.8, 6.7, 4, 0.5, sz=9, c=GREY_LT)

# Right side — placeholder for image
rrect(s, 7.5, 0.5, 5.0, 6.5, GREY_BG)
txt(s, "[ Box illustration ]", 7.5, 3.2, 5.0, 0.5, sz=14, c=GREY_LT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════
# SLIDE 2: HOW IT WORKS
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)

txt(s, "Concept 02: Level Up Series", 0.8, 0.4, 6, 0.3, sz=14, c=BLACK, bold=True)
txt(s, "Iterated Concepts", 8.5, 0.4, 4, 0.3, sz=11, c=GREY_LT, align=PP_ALIGN.RIGHT)
txt(s, "For hypothesis, please refer to Appendix A.", 0.8, 0.75, 6, 0.25, sz=9, c=GREY_LT)

# Divider
rect(s, 0.8, 1.05, 11.7, 0.015, GREY_BG)

# LEFT COLUMN: How it works
txt(s, "How it works", 0.8, 1.2, 5, 0.4, sz=16, c=BLACK, bold=True)

how_items = [
    "Pick up the box at an HPB mall roadshow. Open it that night. Play 5 cards over dinner. That's it \u2014 you've started.",
    "The base box has 4 things inside: conversation cards to break the ice, challenge cards to try something together, a mood check-in QR code, and a quick start guide.",
    "Want to go further? Two add-ons are available at the same roadshow: Wellness Boards (magnetic habit trackers, for P6\u2013Sec 1 only) and a Parent Zone QR card with tips from real parents.",
    "The box isn't homework. Use it once, use it weekly, use it whenever things feel stuck. There's no schedule and no wrong way to play.",
    "Wellness Boards (add-on) go in the youth's room \u2014 not the living room. It's theirs. Parents join only if invited. During exams, swap in the Exam board and focus on sleep.",
    "Parent Zone (add-on) is one QR card. Scan it, watch a 2-minute video from a real parent dealing with the same thing you are. Get one tip. That's it \u2014 not a lecture.",
]
bullets(s, ["\u2022  " + i for i in how_items], 0.8, 1.7, 6.0, 4.5, sz=11, c=DARK, sp=8)

# RIGHT COLUMN: Modality + Design Considerations
txt(s, "Modality", 7.5, 1.2, 5, 0.3, sz=12, c=BLACK, bold=True)
txt(s, "Physical box picked up at HPB mall roadshows. Wellness Boards and Parent Zone QR available as add-ons at the same booth. No app to download. No account to create. QR codes link to digital content when needed.", 7.5, 1.55, 5.0, 1.5, sz=10, c=GREY)

txt(s, "Design Considerations", 7.5, 3.0, 5, 0.3, sz=12, c=BLACK, bold=True)

considerations = [
    "Routine Anchor", "Small Wins", "Autonomy",
    "Physical Prompt", "Gradual Expansion", "Visible Progress",
    "Youth-Owned", "Equal Footing", "Trust First",
    "Transition-Timed", "Replace Don't Restrict",
]
# Tags as a wrapped row
for i, con in enumerate(considerations):
    row = i // 4
    col = i % 4
    cx = 7.5 + col * 1.4
    cy = 3.4 + row * 0.35
    tag(s, con, cx, cy, w=1.3, c=VIOLET)

txt(s, "Base box vs Add-ons", 7.5, 4.7, 5, 0.3, sz=12, c=BLACK, bold=True)
structure = [
    "BASE: Start the Talk, Start Simple, How Are You?, Quick Start Guide",
    "ADD-ON: Wellness Boards (C1) \u2014 5 magnetic habit boards",
    "ADD-ON: Parent Zone (C4) \u2014 Made for Parents QR card",
]
bullets(s, ["\u2022  " + i for i in structure], 7.5, 5.05, 5.0, 1.5, sz=10, c=GREY, sp=4)


# ═══════════════════════════════════════
# SLIDE 3: TOOLS WORKING TOGETHER
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)

txt(s, "Concept 02: Level Up Series - tools working together", 0.8, 0.4, 10, 0.3, sz=14, c=BLACK, bold=True)
txt(s, "Iterated Concepts", 8.5, 0.4, 4, 0.3, sz=11, c=GREY_LT, align=PP_ALIGN.RIGHT)
rect(s, 0.8, 0.75, 11.7, 0.015, GREY_BG)

txt(s, "4 core tools in the base box. 2 optional add-ons available separately. Each tool hands off to the next.", 0.8, 0.85, 10, 0.3, sz=11, c=GREY)

tools = [
    {
        "name": "Start the Talk",
        "color": YELLOW,
        "tier": "CORE",
        "desc": '32 conversation cards. "Never Have I Ever" style. Youth draws first. Both answer honestly. Hit a Swap Card? Answer as each other \u2014 that\'s where it gets fun.',
        "quote": '"Probably start from dinner section first. Then from there can set the same rules to other settings."\nCalvin, parent \u2014 Co-creation',
    },
    {
        "name": "Start Simple",
        "color": CORAL,
        "tier": "CORE",
        "desc": "16 challenge cards + a magnetic scoreboard. Pick a card, do the challenge for a week. Parent plays too \u2014 nobody gets to nag about something they're also being scored on.",
        "quote": '"I\'m guilty of using my headphone while he\'s going to sleep. So this gives me more structure."\nSureena, parent of 13yo, M \u2014 Validation',
    },
    {
        "name": "How Are You?",
        "color": GREEN,
        "tier": "CORE",
        "desc": "One QR card. Scan it, pick how you're feeling, get one small thing to try. No login. No account. Nobody else sees it. Just for you.",
        "quote": '"I wouldn\'t want my parents knowing..."\nChew Xi En, 15 \u2014 Discovery',
    },
    {
        "name": "Quick Start Guide",
        "color": RGBColor(0x26, 0x25, 0x22),
        "tier": "CORE",
        "desc": 'A folded card. Front says "Start here." Back gives you 3 things to try tonight. Inside has the full overview. No prep needed \u2014 just open the box.',
        "quote": "",
    },
    {
        "name": "Wellness Boards",
        "color": BLUE,
        "tier": "ADD-ON",
        "desc": "For P6\u2013Sec 1 only. 5 magnetic boards for the bedroom wall: Sleep, Eat, Move, Mind, and Exam Season. Pick one board, track 2\u20134 habits daily. Start small. During exams, swap in the Exam board and just protect sleep.",
        "quote": '"It\'s more of my thing cuz like I can plan my own schedules."\nEvangeline, 13 \u2014 Validation',
    },
    {
        "name": "Parent Zone",
        "color": VIOLET,
        "tier": "ADD-ON",
        "desc": "One QR card for parents. Scan it to watch a 2-minute video from a real parent dealing with the same thing \u2014 bedtime battles, silent dinners, exam stress. One tip per week via WhatsApp. Not a lecture.",
        "quote": '"I try to ask how school was, but she just says fine. The conversations stopped somewhere around Sec 2."\nP09, parent of 15yo, F \u2014 Discovery',
    },
]

for i, tool in enumerate(tools):
    row = i // 3
    col = i % 3
    cx = 0.8 + col * 4.0
    cy = 1.2 + row * 3.0

    # Card
    rrect(s, cx, cy, 3.75, 2.75, GREY_BG)
    rect(s, cx, cy, 3.75, 0.06, tool["color"])

    # Name
    circ(s, cx + 0.12, cy + 0.2, 0.35, tool["color"])
    txt(s, str(i + 1), cx + 0.12, cy + 0.2, 0.35, 0.35, sz=12, c=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, tool["name"], cx + 0.6, cy + 0.2, 2.8, 0.35, sz=14, c=BLACK, bold=True)

    # Desc
    txt(s, tool["desc"], cx + 0.15, cy + 0.65, 3.45, 1.2, sz=9, c=DARK)

    # Tier badge
    tier_color = RGBColor(0x26, 0x25, 0x22) if tool["tier"] == "CORE" else CORAL
    tag(s, tool["tier"], cx + 2.4, cy + 0.22, w=0.8, c=tier_color)

    # Quote
    if tool["quote"]:
        txt(s, tool["quote"], cx + 0.15, cy + 1.85, 3.45, 0.8, sz=8, c=GREY_LT)


# ═══════════════════════════════════════
# SLIDE 4: WHO IT'S FOR (TRANSITIONS)
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)

txt(s, "Concept 02: Level Up Series - who it's for", 0.8, 0.4, 10, 0.3, sz=14, c=BLACK, bold=True)
txt(s, "Iterated Concepts", 8.5, 0.4, 4, 0.3, sz=11, c=GREY_LT, align=PP_ALIGN.RIGHT)
rect(s, 0.8, 0.75, 11.7, 0.015, GREY_BG)

transitions = [
    {
        "tag": "T1", "period": "P6 \u2192 Sec 1 (12\u201313)", "role": "Play together",
        "color": CORAL,
        "items": [
            "At this age, parents and kids still talk openly. The card game works because it's silly enough to not feel forced.",
            "Start the Talk over dinner. Start Simple as a challenge you both do. Nobody nags \u2014 both play.",
            "Wellness Board add-on is designed for this age. Set it up together \u2014 but it goes in the youth's room, not the living room.",
            "Parent Zone QR helps parents learn how to support without controlling.",
            "The Quick Start Guide tells you what to do tonight. No prep.",
        ],
        "quote": '"Probably start from dinner section first. Then from there can set the same rules to other settings."\nCalvin, parent \u2014 Co-creation',
    },
    {
        "tag": "T2", "period": "Sec 2 \u2192 Sec 3 (14\u201315)", "role": "Friends take over",
        "color": BLUE,
        "items": [
            "By Sec 2, youth don't want parents leading the activity. Friends become the accountability partner instead.",
            "Swap Cards are key here \u2014 answering as each other is less awkward than talking about yourself directly.",
            "Challenge cards shift to friend pairs. The Digital Challenge (app) becomes the main engagement tool.",
            "Wellness Board is fully theirs. Parents step back.",
            "Parent Zone tips shift to: how to stay connected without being annoying.",
        ],
        "quote": '"I try to ask how school was, but she just says fine."\nP09, parent of 15yo, F \u2014 Discovery',
    },
    {
        "tag": "T3", "period": "Sec 4 \u2192 JC/IHL (16+)", "role": "On their own",
        "color": GREEN,
        "items": [
            "Card games might not be needed anymore \u2014 the relationship tools have done their job.",
            "What stays: How Are You? QR. A self-check tool they use on their own.",
            "Wellness Board was for P6\u2013Sec 1. By now they've internalised the habits or moved on.",
            "Parent Zone becomes: I'm here if you need me. Not checking up.",
            "The box started a conversation. What remains is the habit.",
        ],
        "quote": '"Once you reach a certain age you should be more independent... I\'m not going to go and check what they eat."\nP01, parent of 16yo, M \u2014 Discovery',
    },
]

for i, t in enumerate(transitions):
    cx = 0.8 + i * 4.0
    cy = 1.0

    # Card
    rrect(s, cx, cy, 3.75, 6.0, GREY_BG)
    rect(s, cx, cy, 3.75, 0.06, t["color"])

    # Tag
    circ(s, cx + 0.12, cy + 0.2, 0.45, t["color"])
    txt(s, t["tag"], cx + 0.12, cy + 0.2, 0.45, 0.45, sz=14, c=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Period + role
    txt(s, t["period"], cx + 0.7, cy + 0.2, 2.8, 0.25, sz=10, c=t["color"], bold=True)
    txt(s, t["role"], cx + 0.7, cy + 0.45, 2.8, 0.3, sz=16, c=BLACK, bold=True)

    # Bullets
    bullets(s, ["\u2022  " + item for item in t["items"]], cx + 0.15, cy + 0.9, 3.45, 3.5, sz=9, c=DARK, sp=6)

    # Quote at bottom
    txt(s, t["quote"], cx + 0.15, cy + 4.6, 3.45, 1.2, sz=8, c=GREY_LT)


# ═══════════════════════════════════════
# SLIDE 5: BUILDING SUSTAINED USE
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)

txt(s, "Concept 02: Level Up Series - building sustained use", 0.8, 0.4, 10, 0.3, sz=14, c=BLACK, bold=True)
txt(s, "Iterated Concepts", 8.5, 0.4, 4, 0.3, sz=11, c=GREY_LT, align=PP_ALIGN.RIGHT)
txt(s, "For hypothesis, please refer to Appendix A.", 0.8, 0.75, 6, 0.25, sz=9, c=GREY_LT)
rect(s, 0.8, 1.0, 11.7, 0.015, GREY_BG)

# LEFT: Barriers + What would help
txt(s, "Building sustained use", 0.8, 1.15, 5, 0.4, sz=16, c=BLACK, bold=True)

txt(s, "Adoption Barriers", 0.8, 1.7, 5, 0.3, sz=13, c=CORAL, bold=True)
barriers = [
    "The novelty wears off after about 3 weeks. If there's nothing pulling them back, it ends up in a drawer.",
    "If parents start using it to check on them, it feels like monitoring \u2014 and they stop.",
    "During exams, everything healthy gets dropped instead of scaled down.",
    "One bad argument and parents go back to nagging mode. The box can't recover from that alone.",
    "Challenge cards only work if someone else is playing too. If the other person quits, so do they.",
    "The add-ons (Wellness Board, Parent Zone) need to be easy to grab at the roadshow \u2014 if they have to come back for it, they won't.",
]
bullets(s, ["\u2022  " + b for b in barriers], 0.8, 2.05, 5.5, 3.0, sz=10, c=DARK, sp=5)

txt(s, "What would help", 0.8, 4.6, 5, 0.3, sz=13, c=GREEN, bold=True)
helps = [
    "The board goes in the youth's room, not the family area. It never feels like surveillance.",
    "Parents only join if the youth asks them to. That's the rule.",
    "When one board gets boring, switch to the next one. Movement keeps it alive.",
    "During exams, don't quit \u2014 just shrink to one habit. Sleep is the one worth protecting.",
    "The scoreboard means both sides play. Nobody nags about something they're also being scored on.",
    "Parent Zone tips change as the kid gets older \u2014 what works for Sec 1 is different from Sec 3.",
    "How Are You? QR is always there. No schedule, no login, no one watching. Scan it whenever.",
    "At the roadshow, put the add-ons right next to the base box. Don't make people come back for them.",
    "Let youth try the Wellness Board at the booth first \u2014 if they like it, they take one home.",
]
bullets(s, ["\u2022  " + h for h in helps], 0.8, 4.95, 5.5, 2.5, sz=10, c=DARK, sp=4)

# RIGHT: Who it's for
txt(s, "Who it is for", 7.0, 1.15, 5, 0.4, sz=16, c=BLACK, bold=True)

whos = [
    {
        "title": "Sec 1 (12\u201313) \u2014 still open to playing with parents",
        "color": CORAL,
        "items": [
            "The card game works because it's fun, not because it's healthy",
            "Parents can still join without it feeling weird",
            "Start the Talk at dinner \u2192 Start Simple challenge together",
            "Wellness Board add-on: designed for this age group. Set up together, lives in youth's room.",
        ],
    },
    {
        "title": "Sec 2\u20133 (14\u201315) \u2014 friends matter more than parents",
        "color": BLUE,
        "items": [
            "Challenge cards shift to friends, not parents",
            "Swap Cards are less awkward than direct questions at this age",
            "The Digital Challenge app becomes the main tool",
            "If they still have the Wellness Board from Sec 1, it's fully theirs now. Parents back off.",
        ],
    },
    {
        "title": "Sec 4+ (16+) \u2014 managing on their own",
        "color": GREEN,
        "items": [
            "Card games may not be needed anymore",
            "How Are You? QR is what stays \u2014 private, always available",
            "Wellness Board was for younger years. Habits should be internalised by now.",
            "Parents: I'm here if you need me. That's enough.",
        ],
    },
]

curY = 1.7
for who in whos:
    rrect(s, 7.0, curY, 5.5, 1.7, GREY_BG)
    rect(s, 7.0, curY, 0.06, 1.7, who["color"])
    txt(s, who["title"], 7.2, curY + 0.1, 5, 0.3, sz=12, c=who["color"], bold=True)
    bullets(s, ["\u2022  " + i for i in who["items"]], 7.2, curY + 0.4, 5.1, 1.2, sz=9, c=DARK, sp=3)
    curY += 1.85


# ── SAVE ──
out_path = "/Users/alveriacher/Library/CloudStorage/GoogleDrive-alveria.cher@aleph-labs.com/My Drive/HPB YPS/All info/YPSSSSS_Files/Concept 02-  Level Up Series/Concept_02_Level_Up_Series_Updated.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Total: {len(prs.slides)} slides")
